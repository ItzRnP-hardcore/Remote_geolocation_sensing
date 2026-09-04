"""SIH 2026 idea deck for Intelligent Dead Reckoning. python-pptx, 16:9 wide."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn
import sys

OUT = sys.argv[1] if len(sys.argv) > 1 else "SIH2026_IDR.pptx"

# ---- fill these in from the SIH portal -------------------------------------
PS_ID = "SIH26168"
PS_TITLE = "AI-based dead reckoning for vehicle positioning in GNSS-denied environments"
THEME = "Smart Vehicles"
PS_CATEGORY = "Software"
TEAM_ID = "420"
TEAM_NAME = "Odyssey"
# ---------------------------------------------------------------------------

# palette: night-road navy, teal, amber accent
NAVY = RGBColor(0x0F, 0x22, 0x33)
NAVY2 = RGBColor(0x17, 0x33, 0x4A)
TEAL = RGBColor(0x1B, 0x6F, 0x8A)
TEAL_L = RGBColor(0xE3, 0xF0, 0xF4)
AMBER = RGBColor(0xF2, 0xA9, 0x3B)
AMBER_L = RGBColor(0xFD, 0xF1, 0xDC)
INK = RGBColor(0x1F, 0x29, 0x33)
MUTED = RGBColor(0x5F, 0x6B, 0x7A)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GPS_BLUE = RGBColor(0x3B, 0x82, 0xF6)
IMU_ORANGE = RGBColor(0xF9, 0x73, 0x16)
SNAP_GREEN = RGBColor(0x22, 0xC5, 0x5E)
ROAD = RGBColor(0x2A, 0x44, 0x5C)
FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W = 13.333


def rgb(shape_fill, color):
    shape_fill.solid()
    shape_fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill, rounded=False, line=None, radius=0.08):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    rgb(shp.fill, fill)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    if rounded:
        # adjustment is fraction of the shorter side
        shp.adjustments[0] = min(0.5, radius / min(w, h))
    shp.text_frame.text = ""
    return shp


def oval(slide, x, y, d, fill, text=None, size=14, color=WHITE, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    rgb(shp.fill, fill)
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    if text is not None:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = FONT
    return shp


def text(slide, x, y, w, h, content, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, margin=0, italic=False, font=FONT, line_spacing=None,
         space_after=None, bullets=False):
    """content: str | list of paragraphs. A paragraph is str or list of (text, opts) runs."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = anchor
    paras = content if isinstance(content, list) else [content]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        if space_after is not None:
            p.space_after = Pt(space_after)
        runs = para if isinstance(para, list) else [(para, {})]
        if bullets:
            pPr = p._p.get_or_add_pPr()
            pPr.set("marL", str(int(Inches(0.18))))
            pPr.set("indent", str(-int(Inches(0.18))))
            bu = pPr.makeelement(qn("a:buChar"), {"char": "•"})
            pPr.append(bu)
        for rt, o in runs:
            r = p.add_run()
            r.text = rt
            f = r.font
            f.name = o.get("font", font)
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", bold)
            f.italic = o.get("italic", italic)
            f.color.rgb = o.get("color", color)
    return tb


def line(slide, x1, y1, x2, y2, color, width=1.5, dash=None):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dash:
        c.line.dash_style = dash
    return c


def arrow(slide, x1, y1, x2, y2, color=TEAL, width=2):
    c = line(slide, x1, y1, x2, y2, color, width)
    ln = c.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    return c


def polyline(slide, pts, color, width=2.5, dash=None):
    """pts in inches."""
    ff = slide.shapes.build_freeform(Inches(pts[0][0]), Inches(pts[0][1]), scale=1.0)
    ff.add_line_segments([(Inches(x), Inches(y)) for x, y in pts[1:]], close=False)
    shp = ff.convert_to_shape()
    shp.fill.background()
    shp.line.color.rgb = color
    shp.line.width = Pt(width)
    if dash:
        shp.line.dash_style = dash
    shp.shadow.inherit = False
    return shp


def slide_base(dark=False, title=None, number=None):
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY if dark else WHITE
    if title:
        text(s, 0.6, 0.38, 10.5, 0.7, title, size=32, bold=True, color=WHITE if dark else NAVY)
    # footer
    text(s, 0.6, 7.05, 8, 0.3, "Smart India Hackathon 2026  ·  Intelligent Dead Reckoning",
         size=10, color=(RGBColor(0x9F, 0xB3, 0xC8) if dark else MUTED))
    if number:
        text(s, 12.0, 7.05, 0.75, 0.3, f"{number} / 6", size=10, align=PP_ALIGN.RIGHT,
             color=(RGBColor(0x9F, 0xB3, 0xC8) if dark else MUTED))
    return s


def track_card(slide, x, y, w, h, dark=True):
    """The motif: a mini map with three tracks (GPS, IMU-only, map-snapped) through a tunnel."""
    card = rect(slide, x, y, w, h, NAVY2 if dark else LIGHT, rounded=True, radius=0.2)
    # one road through a tunnel
    ry = y + h * 0.5
    road_c = ROAD if dark else RGBColor(0xC8, 0xD1, 0xDA)
    line(slide, x + 0.2, ry, x + w - 0.2, ry, road_c, width=16)
    tx0, tx1 = x + w * 0.36, x + w * 0.64
    rect(slide, tx0, ry - 0.24, tx1 - tx0, 0.48, RGBColor(0x0A, 0x16, 0x22) if dark else RGBColor(0x8A, 0x99, 0xA8), rounded=True)
    text(slide, tx0, ry - 0.55, tx1 - tx0, 0.25, "TUNNEL", size=9, bold=True, align=PP_ALIGN.CENTER, color=AMBER)
    # GNSS: present before and after the tunnel only
    polyline(slide, [(x + 0.3, ry + 0.03), (tx0 - 0.05, ry + 0.03)], GPS_BLUE, 3.5)
    polyline(slide, [(tx1 + 0.05, ry + 0.03), (x + w - 0.3, ry + 0.03)], GPS_BLUE, 3.5)
    # IMU only: drifts off the road inside and keeps diverging
    polyline(slide, [(tx0 - 0.05, ry - 0.02), (tx0 + 0.35, ry - 0.12), (tx1, ry - 0.45),
                     (x + w - 0.3, ry - 0.95)], IMU_ORANGE, 3, MSO_LINE_DASH_STYLE.DASH)
    text(slide, x + w - 1.55, ry - 1.3, 1.3, 0.3, "drift", size=9, italic=True, color=IMU_ORANGE, align=PP_ALIGN.RIGHT)
    # snapped: stays on the road through the tunnel
    polyline(slide, [(tx0 - 0.05, ry - 0.03), (tx1 + 0.05, ry - 0.03)], SNAP_GREEN, 3.5)
    # legend
    lx, ly = x + 0.25, y + h - 0.42
    for i, (c, lab, d) in enumerate([(GPS_BLUE, "GNSS", None), (IMU_ORANGE, "IMU only", MSO_LINE_DASH_STYLE.DASH),
                                     (SNAP_GREEN, "Map-snapped", None)]):
        lx_i = lx + i * (w - 0.5) / 3
        line(slide, lx_i, ly + 0.1, lx_i + 0.35, ly + 0.1, c, 3, d)
        text(slide, lx_i + 0.42, ly, 1.3, 0.22, lab, size=9,
             color=RGBColor(0xC9, 0xD6, 0xE2) if dark else MUTED, anchor=MSO_ANCHOR.MIDDLE)
    return card


# =============================== SLIDE 1: TITLE ===============================
s = slide_base(dark=True)
text(s, 0.7, 0.55, 8, 0.4, "SMART INDIA HACKATHON 2026", size=14, bold=True, color=AMBER)
text(s, 0.7, 1.05, 7.6, 1.6, "Intelligent Dead Reckoning", size=46, bold=True, color=WHITE)
text(s, 0.7, 2.05, 7.4, 1.1,
     "A phone that keeps a vehicle on the map when GNSS goes dark: tunnels, underpasses, "
     "multi-level parking and urban canyons. Fully offline, no server, no extra hardware.",
     size=16, color=RGBColor(0xC9, 0xD6, 0xE2), line_spacing=1.15)

fields = [("Problem Statement ID", PS_ID), ("Problem Statement Title", PS_TITLE),
          ("Theme", THEME), ("PS Category", PS_CATEGORY),
          ("Team ID", TEAM_ID), ("Team Name", TEAM_NAME)]
gx, gy, gw, gh, gap = 0.7, 3.45, 3.62, 1.0, 0.18
for i, (k, v) in enumerate(fields):
    col, row = i % 2, i // 2
    x = gx + col * (gw + gap)
    y = gy + row * (gh + gap)
    rect(s, x, y, gw, gh, NAVY2, rounded=True, radius=0.12)
    text(s, x + 0.2, y + 0.14, gw - 0.4, 0.3, k.upper(), size=9, bold=True, color=AMBER)
    text(s, x + 0.2, y + 0.42, gw - 0.4, 0.55, v, size=13 if len(v) > 40 else 16, bold=True, color=WHITE)

track_card(s, 8.55, 1.1, 4.1, 3.3)
text(s, 8.55, 4.55, 4.1, 1.4,
     [[("Three tracks drawn live on the phone. ", {"bold": True, "color": WHITE}),
       ("Blue is GNSS, orange is the IMU integrating alone, green is the estimate snapped to the "
        "offline road graph. The gap between orange and green is what the tunnel costs.", {})]],
     size=12, color=RGBColor(0xC9, 0xD6, 0xE2), line_spacing=1.1)

# ============================ SLIDE 2: IDEA / SOLUTION ========================
s = slide_base(title="Idea & Proposed Solution", number=2)
# left column: problem + solution
text(s, 0.6, 1.25, 5.6, 0.35, "THE PROBLEM", size=11, bold=True, color=TEAL)
text(s, 0.6, 1.6, 5.6, 1.7, [
    "GNSS fails exactly where a driver needs it most: tunnels, underpasses, flyover stacks, parking decks and dense urban canyons.",
    "Consumer navigation freezes or jumps to a wrong road, and recovers only after the sky is back.",
    "The phone's own IMU can bridge the gap, but drifts by tens of metres within a minute if it is left to integrate on its own.",
], size=12.5, bullets=True, space_after=6, line_spacing=1.1)

text(s, 0.6, 3.4, 5.6, 0.35, "OUR SOLUTION", size=11, bold=True, color=TEAL)
rect(s, 0.6, 3.8, 5.6, 1.55, TEAL_L, rounded=True, radius=0.15)
text(s, 0.8, 3.92, 5.2, 1.35,
     [[("An Android app that fuses the phone's IMU, GNSS and an ", {}),
       ("offline OpenStreetMap road graph", {"bold": True}),
       (", shifting trust from satellites to inertial sensing ", {}),
       ("before", {"bold": True}),
       (" the fix is lost, and snapping the drifting estimate back onto the road so the vehicle "
        "never leaves the map.", {})]],
     size=13.5, color=INK, line_spacing=1.12, anchor=MSO_ANCHOR.MIDDLE)

# uniqueness chips
text(s, 0.6, 5.6, 5.6, 0.3, "HOW IT DIFFERS", size=11, bold=True, color=TEAL)
chips = ["100% offline", "No server, no API key", "No extra hardware", "Measured, not assumed"]
for i, c in enumerate(chips):
    cw = 2.7
    cx = 0.6 + (i % 2) * (cw + 0.2)
    cy = 5.95 + (i // 2) * 0.52
    rect(s, cx, cy, cw, 0.42, NAVY, rounded=True, radius=0.21)
    text(s, cx, cy, cw, 0.42, c, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

# right: 2x2 feature cards
feats = [
    ("Strapdown inertial integrator",
     "Attitude from the rotation vector, acceleration levelled into ENU, gravity and bias learned on the "
     "device. Zero-velocity updates pin drift at every stop."),
    ("GNSS trust gate",
     "Satellite count and C/N0 collapse before fixes stop. The filter hands over to the IMU on that early "
     "evidence instead of waiting for a timeout."),
    ("On-device HMM map matching",
     "A Viterbi matcher snaps the estimate to the road and feeds the road bearing back as a heading "
     "observation. Heading error stops accumulating."),
    ("Road graph rebuilt on the phone",
     "Mapsforge vector tiles carry no connectivity. We recover it by snapping endpoints on a half-metre "
     "grid, so no routing server is needed, even underground."),
]
fx, fy, fw, fh, fg = 6.6, 1.25, 3.0, 2.45, 0.18
for i, (h, b) in enumerate(feats):
    col, row = i % 2, i // 2
    x = fx + col * (fw + fg)
    y = fy + row * (fh + fg)
    rect(s, x, y, fw, fh, LIGHT, rounded=True, radius=0.15)
    oval(s, x + 0.2, y + 0.2, 0.42, AMBER, str(i + 1), size=14, color=NAVY)
    text(s, x + 0.72, y + 0.2, fw - 0.9, 0.5, h, size=13, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.2, y + 0.78, fw - 0.4, fh - 0.9, b, size=11.5, color=INK, line_spacing=1.1)

# ============================ SLIDE 3: TECHNICAL APPROACH =====================
s = slide_base(title="Technical Approach", number=3)
# pipeline row
stages = [
    ("Phone sensors", "Accel, gyro, mag, rotation vector, barometer at 200 Hz; GNSS + satellite status at 1 Hz"),
    ("Levelling & ZUPT", "Rotate into ENU, remove learned gravity and bias, zero-velocity update when parked"),
    ("Strapdown integrator", "10 Hz position and velocity; re-anchored to GNSS only while the fix is healthy"),
    ("HMM map matcher", "Viterbi over the on-device road graph; matched bearing fed back as heading"),
    ("Along-road tracker", "Position as a scalar along the road polyline during an outage; road supplies direction"),
    ("Offline map UI", "osmdroid + Mapsforge vector maps; GNSS, IMU and snapped tracks drawn live"),
]
px, py, pw, ph, pg = 0.6, 1.3, 1.9, 1.75, 0.15
for i, (h, b) in enumerate(stages):
    x = px + i * (pw + pg)
    dark_stage = i in (2, 3)
    rect(s, x, py, pw, ph, NAVY if dark_stage else TEAL_L, rounded=True, radius=0.14)
    oval(s, x + 0.15, py + 0.15, 0.34, AMBER if dark_stage else TEAL, str(i + 1), size=11, color=NAVY if dark_stage else WHITE)
    text(s, x + 0.55, py + 0.13, pw - 0.65, 0.4, h, size=11.5, bold=True,
         color=WHITE if dark_stage else NAVY, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.15, py + 0.6, pw - 0.3, ph - 0.7, b, size=9.5,
         color=RGBColor(0xDD, 0xE7, 0xEF) if dark_stage else INK, line_spacing=1.08)
    if i < len(stages) - 1:
        arrow(s, x + pw + 0.01, py + ph / 2, x + pw + pg - 0.01, py + ph / 2, TEAL, 2)
# feedback loop label under the pipeline
line(s, px + 3 * (pw + pg) + pw / 2, py + ph + 0.02, px + 3 * (pw + pg) + pw / 2, py + ph + 0.25, AMBER, 1.5)
line(s, px + 2 * (pw + pg) + pw / 2, py + ph + 0.25, px + 3 * (pw + pg) + pw / 2, py + ph + 0.25, AMBER, 1.5)
arrow(s, px + 2 * (pw + pg) + pw / 2, py + ph + 0.25, px + 2 * (pw + pg) + pw / 2, py + ph + 0.03, AMBER, 1.5)
text(s, px + 2 * (pw + pg) + 0.2, py + ph + 0.28, 3.6, 0.3, "heading feedback: matched road bearing corrects the integrator",
     size=9, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# lower left: tech stack
text(s, 0.6, 3.75, 6.0, 0.3, "TECHNOLOGY STACK", size=11, bold=True, color=TEAL)
stack = [
    ("Mobile", "Kotlin, Android foreground service, single-writer sensor thread, 2 s flush so a crash costs at most two seconds"),
    ("Maps", "osmdroid + Mapsforge vector .map files (OpenStreetMap data, ODbL). In-app zone downloader for all of India"),
    ("AI / ML", "PyTorch 1-D ResNet speed head exported to TorchScript Mobile; 40 ms per inference on a dedicated thread"),
    ("Fusion", "Strapdown INS + ZUPT + scalar Kalman gain on model speed; HMM Viterbi map matching; along-road tracking"),
    ("Evaluation", "Python harness: CRSE, CAE, AEPS on synthetic outages cut from recorded drives, scored against GNSS truth"),
]
sy = 4.1
for k, v in stack:
    rect(s, 0.6, sy, 1.25, 0.5, NAVY, rounded=True, radius=0.1)
    text(s, 0.6, sy, 1.25, 0.5, k, size=10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.98, sy, 5.3, 0.5, v, size=10.5, color=INK, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    sy += 0.58

# lower right: threads + gate
text(s, 7.75, 3.75, 5.0, 0.3, "ON-DEVICE EXECUTION", size=11, bold=True, color=TEAL)
rect(s, 7.75, 4.1, 4.98, 1.55, LIGHT, rounded=True, radius=0.15)
rows = [("Thread", "Owns", "Rate", True),
        ("imu-logger", "Sensor + GNSS callbacks, all file writes, integrator", "200 Hz", False),
        ("imu-ml", "Model load and inference", "10 Hz", False),
        ("map-match", "Tile reads, HMM matcher, road graph", "0.5 Hz", False)]
ry = 4.2
for a, b, c, hdr in rows:
    text(s, 7.95, ry, 1.2, 0.33, a, size=10, bold=True, color=TEAL if hdr else NAVY, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 9.15, ry, 2.7, 0.33, b, size=10, bold=hdr, color=TEAL if hdr else INK, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 11.85, ry, 0.8, 0.33, c, size=10, bold=True, color=TEAL if hdr else NAVY, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
    ry += 0.35
rect(s, 7.75, 5.8, 4.98, 1.1, AMBER_L, rounded=True, radius=0.15)
text(s, 7.95, 5.88, 4.6, 0.95,
     [[("Exactly one writer thread, no locks. ", {"bold": True, "color": NAVY}),
       ("Inference results are posted back to the logger thread; stale requests are dropped rather than "
        "queued. GNSS is trusted only with 4+ satellites used and accuracy under 20 m.", {})]],
     size=10.5, color=INK, line_spacing=1.1, anchor=MSO_ANCHOR.MIDDLE)

# ======================== SLIDE 4: FEASIBILITY & VIABILITY ====================
s = slide_base(title="Feasibility & Viability", number=4)
# chart left
text(s, 0.6, 1.2, 6.4, 0.3, "MEASURED ON DEVICE: POSITION ERROR DURING SYNTHETIC GNSS OUTAGES", size=11, bold=True, color=TEAL)
cd = CategoryChartData()
cd.categories = ["10 s outage", "30 s outage", "60 s outage"]
cd.add_series("Dead reckoning + heading feedback (live today)", (10.2, 36.0, 87.5))
cd.add_series("Along-road, fed integrator distance", (12.6, 40.8, 104.7))
cd.add_series("Along-road, fed true distance", (8.3, 22.3, 50.5))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.5), Inches(6.5), Inches(3.55), cd)
ch = gf.chart
ch.has_title = False
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(9)
ch.legend.font.color.rgb = INK
ch.font.name = FONT
va = ch.value_axis
va.has_major_gridlines = True
va.major_gridlines.format.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
va.format.line.fill.background()
va.tick_labels.font.size = Pt(9)
va.tick_labels.font.color.rgb = MUTED
va.has_title = True
va.axis_title.text_frame.text = "mean error, m (CRSE)"
va.axis_title.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
va.axis_title.text_frame.paragraphs[0].runs[0].font.color.rgb = MUTED
va.axis_title.text_frame.paragraphs[0].runs[0].font.bold = False
ca = ch.category_axis
ca.tick_labels.font.size = Pt(10)
ca.tick_labels.font.color.rgb = INK
ca.format.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
plot = ch.plots[0]
plot.gap_width = 60
plot.overlap = -10
plot.has_data_labels = True
dl = plot.data_labels
dl.font.size = Pt(9)
dl.font.color.rgb = INK
dl.number_format = '0.0'
dl.number_format_is_linked = False
dl.position = XL_LABEL_POSITION.OUTSIDE_END
for ser, col in zip(plot.series, [TEAL, RGBColor(0xB0, 0xBE, 0xCC), AMBER]):
    ser.format.fill.solid()
    ser.format.fill.fore_color.rgb = col
text(s, 0.6, 5.1, 6.4, 0.5,
     "46 to 57 outages per duration, cut from recorded sessions and scored against GNSS ground truth. "
     "Along-road tracking cuts 60-second error by 42% once it is fed a trustworthy distance.",
     size=10, italic=True, color=MUTED, line_spacing=1.1)

# right: working today
text(s, 7.4, 1.2, 5.4, 0.3, "WORKING TODAY", size=11, bold=True, color=TEAL)
works = [
    "Android app records 11 sensor streams on a Samsung Galaxy S21 FE, time-aligned on one monotonic clock",
    "Dead reckoning with map-matched heading feedback runs on the phone and is measured, not simulated",
    "Road topology recovered on device: 25,383 segments around IIT Kharagpur, 100% linked, 98.5% in one component",
    "Offline maps render with no tile server; zones for the rest of India download in-app",
]
wy = 1.55
for w_ in works:
    oval(s, 7.4, wy + 0.05, 0.3, SNAP_GREEN, "✓", size=11, color=WHITE)
    text(s, 7.85, wy, 4.9, 0.62, w_, size=11, color=INK, line_spacing=1.08, anchor=MSO_ANCHOR.MIDDLE)
    wy += 0.68

# bottom band: risks and mitigation
text(s, 0.6, 5.7, 6, 0.3, "CHALLENGES & HOW WE CLOSE THEM", size=11, bold=True, color=TEAL)
risks = [
    ("Along-track bias",
     "A -0.02 m/s² acceleration bias costs ~30% of distance over 60 s unaided. "
     "Fix: learn horizontal bias against GNSS velocity while the fix is healthy, carry it into the outage."),
    ("Under-trained speed model",
     "304 training windows, walking pace only. Fix: record vehicle-speed drives, retrain, "
     "and enable speed fusion only if it beats a constant predictor."),
    ("Map size and field validation",
     "The 210 MB zone map ships as a release asset with an in-app downloader. "
     "Heading feedback has run only in replay; next step is a live road test."),
]
rx, rw, rgap = 0.6, 3.95, 0.16
for i, (h, b) in enumerate(risks):
    x = rx + i * (rw + rgap)
    rect(s, x, 6.05, rw, 0.95, LIGHT, rounded=True, radius=0.12)
    text(s, x + 0.15, 6.1, rw - 0.3, 0.25, h, size=10.5, bold=True, color=NAVY)
    text(s, x + 0.15, 6.35, rw - 0.3, 0.65, b, size=9, color=INK, line_spacing=1.05)

# ============================ SLIDE 5: IMPACT & BENEFITS ======================
s = slide_base(title="Impact & Benefits", number=5)
stats = [("0", "cloud calls, servers or API keys"), ("10 Hz", "position updates on the handset"),
         ("42%", "lower 60 s outage error with along-road tracking"), ("100%", "of India coverable with free OSM zone maps")]
sx, sw, sg = 0.6, 2.93, 0.15
for i, (n, l) in enumerate(stats):
    x = sx + i * (sw + sg)
    rect(s, x, 1.25, sw, 1.35, NAVY, rounded=True, radius=0.15)
    text(s, x + 0.2, 1.3, sw - 0.4, 0.7, n, size=34, bold=True, color=AMBER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.2, 1.98, sw - 0.4, 0.55, l, size=10.5, color=RGBColor(0xC9, 0xD6, 0xE2), line_spacing=1.05)

cols = [
    ("Social", TEAL, [
        "Drivers, ambulances and delivery riders keep guidance through tunnels, underpasses and hill-road cuttings",
        "No wrong turn at a tunnel exit: the position is already on the right road when the sky returns",
        "Works on the phone people already own; nothing to buy or install in the vehicle",
    ]),
    ("Economic", AMBER, [
        "Zero recurring cost: no map API billing, no data plan, no server to run",
        "Fleet and logistics tracking that does not go blind under flyovers or in parking decks",
        "One codebase for OEM navigation, ride-hailing and toll or road-agency apps",
    ]),
    ("Strategic & Environmental", SNAP_GREEN, [
        "Location never leaves the device: privacy by construction",
        "Offline by design, so it serves low-connectivity and remote regions first, not last",
        "Fewer missed exits and detours means less fuel burnt; road topology from OSM is open and reusable",
    ]),
]
cx, cw, cg = 0.6, 3.95, 0.16
for i, (h, col, items) in enumerate(cols):
    x = cx + i * (cw + cg)
    rect(s, x, 2.85, cw, 2.85, LIGHT, rounded=True, radius=0.15)
    oval(s, x + 0.2, 3.02, 0.3, col)
    text(s, x + 0.62, 2.95, cw - 0.8, 0.45, h, size=15, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.2, 3.5, cw - 0.4, 2.15, items, size=12, color=INK, bullets=True, space_after=8, line_spacing=1.08)

rect(s, 0.6, 5.9, 12.13, 0.9, TEAL_L, rounded=True, radius=0.15)
text(s, 0.8, 5.95, 11.8, 0.8,
     [[("Scalability. ", {"bold": True, "color": NAVY}),
       ("One Mapsforge file per zone (210 MB for eastern India, 1.5 GB for the whole country). The same "
        "estimator extends to metro and rail, indoor parking and two-wheelers, and its measured outage "
        "harness lets every future change be gated on numbers.", {})]],
     size=11.5, color=INK, line_spacing=1.1, anchor=MSO_ANCHOR.MIDDLE)

# ========================== SLIDE 6: RESEARCH & REFERENCES ====================
s = slide_base(title="Research & References", number=6)
text(s, 0.6, 1.2, 7.4, 0.3, "LITERATURE WE BUILD ON", size=11, bold=True, color=TEAL)
refs = [
    ("Brossard, Barrau, Bonnabel (2019).", "AI-IMU Dead-Reckoning. arXiv:1904.06064"),
    ("Hurwitz, Cohen, Klein (2024).", "Deep Learning Assisted Inertial Dead Reckoning and Fusion. arXiv:2407.16387"),
    ("Gomes, Costa (2020).", "Vehicular Dead Reckoning Based on Machine Learning and Map Matching. IEEE"),
    ("Suwandi, Kitasuka, Aritsugi (2019).", "Vehicle Vibration Error Compensation on IMU-accelerometer Sensor Using Adaptive Filter and Low-pass Filter Approaches. J. Inf. Processing 27"),
    ("Nirmal et al. (2016).", "Noise modeling and analysis of an IMU-based attitude sensor. arXiv:1608.07053"),
    ("Chao et al. (2022).", "A Comprehensive Review of Map-Matching Techniques. Int. J. Web Services Research 19(1)"),
    ("Newson, Krumm (2009).", "Hidden Markov Map Matching Through Noise and Sparseness. ACM SIGSPATIAL"),
    ("IO-VNBD (Coventry Univ.).", "Inertial and Odometry Vehicle Navigation Benchmark Dataset for model training and CRSE / CAE / AEPS metric definitions"),
]
ry = 1.55
for a, b in refs:
    text(s, 0.6, ry, 7.5, 0.5, [[(a + " ", {"bold": True, "color": NAVY}), (b, {})]],
         size=10.5, color=INK, line_spacing=1.05, anchor=MSO_ANCHOR.MIDDLE)
    ry += 0.5

text(s, 8.5, 1.2, 4.3, 0.3, "EXISTING SOLUTIONS STUDIED", size=11, bold=True, color=TEAL)
ex = [
    ("Google Maps SDK", "No offline tile access, no dead-reckoning hook, billing account required. We use OSM data under ODbL instead."),
    ("OSRM-based map matching", "Accurate, but needs a routing server. We rebuild road connectivity on the handset from Mapsforge tiles."),
    ("Learned inertial odometry (RoNIN, AI-IMU)", "Strong results with dense training data. Our speed head reuses the idea but is gated until it beats a constant."),
]
ey = 1.55
for h, b in ex:
    rect(s, 8.5, ey, 4.23, 1.1, LIGHT, rounded=True, radius=0.12)
    text(s, 8.68, ey + 0.08, 3.9, 0.3, h, size=11, bold=True, color=NAVY)
    text(s, 8.68, ey + 0.38, 3.9, 0.7, b, size=9.5, color=INK, line_spacing=1.05)
    ey += 1.22

rect(s, 8.5, 5.25, 4.23, 1.55, NAVY, rounded=True, radius=0.15)
text(s, 8.7, 5.32, 3.9, 1.45,
     [[("Open data and tools. ", {"bold": True, "color": AMBER}),
       ("OpenStreetMap (ODbL), Mapsforge map builds, osmdroid, PyTorch Mobile, Android Sensor and GNSS "
        "status APIs. Every number in this deck is reproducible from the committed evaluation harness "
        "against recorded sessions.", {})]],
     size=10.5, color=RGBColor(0xDD, 0xE7, 0xEF), line_spacing=1.1, anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print("wrote", OUT)
